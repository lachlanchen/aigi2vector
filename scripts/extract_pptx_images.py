#!/usr/bin/env python3
"""Extract embedded images from a PPTX by slide."""

import argparse
import hashlib
import io
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract embedded images from a PPTX by slide."
    )
    parser.add_argument("pptx_path", help="Path to .pptx file")
    parser.add_argument(
        "output_dir",
        help="Directory to write extracted images",
    )
    parser.add_argument(
        "--dedupe",
        action="store_true",
        help="Skip duplicate images (by SHA1 hash)",
    )
    parser.add_argument(
        "--png",
        action="store_true",
        help="Convert and save images as PNG",
    )
    return parser.parse_args()


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def write_image(
    out_dir: Path,
    slide_idx: int,
    image_idx: int,
    image,
    as_png: bool,
) -> Path:
    if as_png:
        filename = f"p{slide_idx}image{image_idx}.png"
        out_path = out_dir / filename
        try:
            with Image.open(io.BytesIO(image.blob)) as im:
                im.save(out_path, format="PNG")
            return out_path
        except Exception:
            # Fall back to original bytes if conversion fails
            pass

    ext = image.ext
    filename = f"p{slide_idx}image{image_idx}.{ext}"
    out_path = out_dir / filename
    with open(out_path, "wb") as f:
        f.write(image.blob)
    return out_path


def main() -> None:
    args = parse_args()
    pptx_path = Path(args.pptx_path)
    out_dir = Path(args.output_dir)

    if not pptx_path.exists():
        raise FileNotFoundError(pptx_path)

    ensure_dir(out_dir)

    prs = Presentation(str(pptx_path))

    seen = set()
    total = 0
    for slide_i, slide in enumerate(prs.slides, start=1):
        image_i = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            image = shape.image
            if args.dedupe:
                digest = hashlib.sha1(image.blob).hexdigest()
                if digest in seen:
                    continue
                seen.add(digest)
            image_i += 1
            write_image(out_dir, slide_i, image_i, image, args.png)
            total += 1

    print(f"Extracted {total} images to {out_dir}")


if __name__ == "__main__":
    main()
